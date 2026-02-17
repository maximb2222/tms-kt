from __future__ import annotations

from pathlib import Path
import re

from pptx import Presentation
from pptx.util import Pt
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas


ROOT = Path(__file__).resolve().parent.parent
DOCS_DIR = ROOT / "docs"


def register_font() -> tuple[str, str]:
    candidates = [
        (Path(r"C:\Windows\Fonts\arial.ttf"), Path(r"C:\Windows\Fonts\arialbd.ttf")),
        (Path(r"C:\Windows\Fonts\calibri.ttf"), Path(r"C:\Windows\Fonts\calibrib.ttf")),
        (Path(r"C:\Windows\Fonts\tahoma.ttf"), Path(r"C:\Windows\Fonts\tahomabd.ttf")),
    ]

    for regular, bold in candidates:
        if regular.exists() and bold.exists():
            pdfmetrics.registerFont(TTFont("DocRegular", str(regular)))
            pdfmetrics.registerFont(TTFont("DocBold", str(bold)))
            return "DocRegular", "DocBold"

    raise FileNotFoundError("No suitable Cyrillic font found in C:\\Windows\\Fonts.")


def normalize_markdown_lines(content: str) -> list[tuple[str, bool]]:
    result: list[tuple[str, bool]] = []
    for raw in content.splitlines():
        line = raw.rstrip()
        if not line.strip():
            result.append(("", False))
            continue

        if line.startswith("#"):
            title = line.lstrip("#").strip()
            result.append((title, True))
            continue

        if line.startswith("|"):
            cells = [cell.strip() for cell in line.strip("|").split("|")]
            line = " | ".join(cells)

        result.append((line, False))

    return result


def wrap_by_width(text: str, font_name: str, font_size: int, max_width: float) -> list[str]:
    if not text:
        return [""]

    words = text.split()
    if not words:
        return [""]

    lines: list[str] = []
    current = words[0]

    for word in words[1:]:
        candidate = f"{current} {word}"
        if pdfmetrics.stringWidth(candidate, font_name, font_size) <= max_width:
            current = candidate
        else:
            lines.append(current)
            current = word

    lines.append(current)
    return lines


def markdown_to_pdf(source_md: Path, target_pdf: Path) -> None:
    regular_font, bold_font = register_font()
    page_width, page_height = A4
    left = 40
    right = 40
    top = 50
    bottom = 45
    max_width = page_width - left - right

    pdf = canvas.Canvas(str(target_pdf), pagesize=A4)
    pdf.setTitle(source_md.stem)
    y = page_height - top

    def ensure_space(line_height: int) -> None:
        nonlocal y
        if y - line_height < bottom:
            pdf.showPage()
            y = page_height - top

    lines = normalize_markdown_lines(source_md.read_text(encoding="utf-8"))
    for line, is_heading in lines:
        if not line:
            y -= 8
            continue

        font_name = bold_font if is_heading else regular_font
        font_size = 13 if is_heading else 11
        line_height = 18 if is_heading else 15

        bullet_match = re.match(r"^(\s*[-*]\s+)(.+)$", line)
        if bullet_match:
            prefix = bullet_match.group(1)
            body = bullet_match.group(2)
            prefix_width = pdfmetrics.stringWidth(prefix, font_name, font_size)
            wrapped = wrap_by_width(body, font_name, font_size, max_width - prefix_width)
            for index, segment in enumerate(wrapped):
                ensure_space(line_height)
                text = f"{prefix}{segment}" if index == 0 else (" " * len(prefix) + segment)
                pdf.setFont(font_name, font_size)
                pdf.drawString(left, y, text)
                y -= line_height
            continue

        wrapped = wrap_by_width(line, font_name, font_size, max_width)
        for segment in wrapped:
            ensure_space(line_height)
            pdf.setFont(font_name, font_size)
            pdf.drawString(left, y, segment)
            y -= line_height

    pdf.save()


def markdown_to_pptx(source_md: Path, target_pptx: Path) -> None:
    content = source_md.read_text(encoding="utf-8").splitlines()
    slides: list[tuple[str, list[str]]] = []
    current_title = "Overview"
    bullets: list[str] = []

    for raw in content:
        line = raw.strip()
        if not line:
            continue

        if line.startswith("# "):
            continue

        if line.startswith("## "):
            if bullets:
                slides.append((current_title, bullets))
            current_title = line[3:].strip()
            bullets = []
            continue

        if line.startswith("- "):
            bullets.append(line[2:].strip())
            continue

        numbered = re.match(r"^\d+\.\s+(.*)$", line)
        if numbered:
            bullets.append(numbered.group(1).strip())
            continue

        bullets.append(line)

    if bullets:
        slides.append((current_title, bullets))

    prs = Presentation()

    for title, slide_bullets in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(30)

        body = slide.shapes.placeholders[1].text_frame
        body.clear()

        for index, bullet in enumerate(slide_bullets[:8]):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(20)
            paragraph.font.name = "Calibri"

    prs.save(target_pptx)


def main() -> None:
    markdown_to_pdf(DOCS_DIR / "KT1-Plan.md", DOCS_DIR / "KT1-Plan.pdf")
    markdown_to_pdf(DOCS_DIR / "KT3-FinalDocumentation.md", DOCS_DIR / "KT3-FinalDocumentation.pdf")
    markdown_to_pptx(DOCS_DIR / "KT4-Presentation-Material.md", DOCS_DIR / "KT4-Presentation-Material.pptx")
    print("Generated: KT1-Plan.pdf, KT3-FinalDocumentation.pdf, KT4-Presentation-Material.pptx")


if __name__ == "__main__":
    main()
